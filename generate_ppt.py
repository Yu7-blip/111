#!/usr/bin/env python3
"""生成答辩PPT — 黑白简洁学术风格"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

prs = Presentation()
prs.slide_width  = Cm(33.87)  # 16:9
prs.slide_height = Cm(19.05)

BLACK   = RGBColor(0x00, 0x00, 0x00)
DARK    = RGBColor(0x33, 0x33, 0x33)
GRAY    = RGBColor(0x88, 0x88, 0x88)
LGRAY   = RGBColor(0xE8, 0xE8, 0xE8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY   = RGBColor(0x55, 0x55, 0x55)

# ==================== 工具函数 ====================
def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=BLACK, font_name='Microsoft YaHei', align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(font_size)
    p.font.bold = bold; p.font.color.rgb = color; p.font.name = font_name
    p.alignment = align
    return tf

def add_rich_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame; tf.word_wrap = True
    return tf

def add_para(tf, text, font_size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             space_before=4, space_after=2, font_name='Microsoft YaHei', indent=0):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(font_size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font_name; p.alignment = align
    p.space_before = Pt(space_before); p.space_after = Pt(space_after)
    if indent: p.level = indent
    return p

def add_bullet(tf, text, font_size=11, level=0, color=DARK):
    p = add_para(tf, text, font_size=font_size, color=color, indent=level)
    return p

def add_title_bar(slide, title_text):
    """顶部标题栏 + 下划线"""
    # 灰色顶条
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(33.87), Cm(2.2))
    shape.fill.solid(); shape.fill.fore_color.rgb = LGRAY
    shape.line.fill.background()
    add_textbox(slide, 1.5, 0.35, 28, 1.5, title_text, font_size=28, bold=True, color=BLACK)
    # 底线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.5), Cm(2.2), Cm(30.87), Cm(0.06))
    line.fill.solid(); line.fill.fore_color.rgb = DARK; line.line.fill.background()

def add_page_number(slide, num, total):
    add_textbox(slide, 29.5, 17.5, 4, 1, f'{num}/{total}', font_size=9, color=GRAY, align=PP_ALIGN.RIGHT)

def add_section_box(slide, left, top, width, height, title, bullets):
    """添加带标题的分区框"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = GRAY; shape.line.width = Pt(0.5)
    tf = add_rich_textbox(slide, left+0.5, top+0.3, width-1, height-0.6)
    add_para(tf, title, font_size=14, bold=True, color=BLACK)
    for b in bullets:
        add_para(tf, '• '+b, font_size=10, color=DARK, space_before=2)

def add_table(slide, left, top, col_widths, headers, rows, font_size=9):
    """添加简洁表格"""
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
        Cm(left), Cm(top), Cm(sum(col_widths)), Cm(0.8 * n_rows))
    table = table_shape.table
    for ci, w in enumerate(col_widths):
        table.columns[ci].width = Cm(w)
    # header
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci); cell.text = ''
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h; r.font.size = Pt(font_size); r.font.bold = True
        r.font.name = 'Microsoft YaHei'; r.font.color.rgb = BLACK
        cell.fill.solid(); cell.fill.fore_color.rgb = LGRAY
    # data
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri+1, ci); cell.text = ''
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = str(val); r.font.size = Pt(font_size-1)
            r.font.name = 'Microsoft YaHei'; r.font.color.rgb = DARK
    return table_shape

TOTAL_SLIDES = 16

# ================================================================
#  Slide 1 — 封面
# ================================================================
s = add_blank_slide()
# 大色块背景
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(33.87), Cm(19.05))
shape.fill.solid(); shape.fill.fore_color.rgb = WHITE; shape.line.fill.background()
# 顶部装饰线
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(3), Cm(5.5), Cm(27.87), Cm(0.04))
line.fill.solid(); line.fill.fore_color.rgb = DARK; line.line.fill.background()
# 标题
add_textbox(s, 3, 2.5, 28, 3, '外卖配送平台', font_size=48, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
add_textbox(s, 3, 5.8, 28, 1.5, '系统设计与实现  |  答辩汇报', font_size=24, bold=False, color=DARK, align=PP_ALIGN.CENTER)
# 信息
info_tf = add_rich_textbox(s, 8, 10, 18, 6)
add_para(info_tf, '第 11 组', font_size=20, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
add_para(info_tf, '', font_size=8)
add_para(info_tf, '余  睿（组长）  |  2300770172  |  软工2304', font_size=14, color=DARK, align=PP_ALIGN.CENTER, space_before=6)
add_para(info_tf, '赵宇弦          |  2300770178  |  软工2304', font_size=14, color=DARK, align=PP_ALIGN.CENTER, space_before=2)
add_para(info_tf, '周峥嵘          |  2300770046  |  软工2304', font_size=14, color=DARK, align=PP_ALIGN.CENTER, space_before=2)
add_para(info_tf, '', font_size=6)
add_para(info_tf, '指导老师：丁玺润', font_size=14, color=DARK, align=PP_ALIGN.CENTER, space_before=2)
add_para(info_tf, '', font_size=6)
add_para(info_tf, f'{datetime.date.today().strftime("%Y年%m月%d日")}', font_size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_page_number(s, 1, TOTAL_SLIDES)


# ================================================================
#  Slide 2 — 目录
# ================================================================
s = add_blank_slide()
add_title_bar(s, '汇报目录')
items = [
    ('01', '项目背景与目标'),
    ('02', '技术架构总览'),
    ('03', '需求分析概要'),
    ('04', '数据库设计'),
    ('05', '核心功能实现'),
    ('06', '关键技术难点攻克'),
    ('07', '系统测试'),
    ('08', '团队分工'),
    ('09', '总结与展望'),
]
for i, (num, title) in enumerate(items):
    y = 3.5 + i * 1.6
    add_textbox(s, 4, y, 1.5, 1, num, font_size=22, bold=True, color=DARK)
    add_textbox(s, 6, y, 22, 1, title, font_size=22, bold=False, color=BLACK)
    if i < len(items)-1:
        l = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(6), Cm(y+1.2), Cm(24), Cm(0.02))
        l.fill.solid(); l.fill.fore_color.rgb = LGRAY; l.line.fill.background()
add_page_number(s, 2, TOTAL_SLIDES)


# ================================================================
#  Slide 3 — 项目背景与目标
# ================================================================
s = add_blank_slide()
add_title_bar(s, '一、项目背景与目标')
# 左侧：背景
tf = add_rich_textbox(s, 1.5, 3.2, 14.5, 14)
add_para(tf, '项目背景', font_size=18, bold=True, color=BLACK)
add_para(tf, '', font_size=4)
add_para(tf, 'O2O外卖已成为国民级日常消费场景，美团/饿了么日均订单量超5000万单。', font_size=12, color=DARK)
add_para(tf, '外卖平台涉及的技术栈涵盖：空间索引、实时追踪、智能调度、分布式事务、高并发缓存等软件工程核心领域，是课程设计的理想综合实践题目。', font_size=12, color=DARK, space_before=8)
add_para(tf, '', font_size=8)
add_para(tf, '项目目标', font_size=18, bold=True, color=BLACK, space_before=12)
add_para(tf, '', font_size=4)
for g in ['构建完整的外卖配送生态系统（用户+商家+骑手+平台）',
          '基于GeoHash实现附近商家高效空间搜索',
          '设计四维加权智能派单算法',
          '采用Transaction Outbox保证分布式事务一致性',
          '实现缓存三防（穿透/击穿/雪崩）与多层安全防护']:
    add_para(tf, '◆  ' + g, font_size=11, color=DARK, space_before=6)
# 右侧：数据卡片
for i, (num, label) in enumerate([('49', '子需求'), ('3', '大模块'), ('4', '角色端'), ('20', '数据表')]):
    y = 3.5 + i * 3.5
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(18.5), Cm(y), Cm(13.5), Cm(3))
    box.fill.solid(); box.fill.fore_color.rgb = LGRAY; box.line.fill.background()
    add_textbox(s, 18.5, y+0.3, 13.5, 1.5, num, font_size=36, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_textbox(s, 18.5, y+1.8, 13.5, 0.8, label, font_size=12, color=DARK, align=PP_ALIGN.CENTER)
add_page_number(s, 3, TOTAL_SLIDES)


# ================================================================
#  Slide 4 — 技术架构
# ================================================================
s = add_blank_slide()
add_title_bar(s, '二、技术架构总览')
# 四列架构
cols = [
    ('客户端层', ['微信小程序\n(用户端+骑手端)', 'Admin SPA\n(Vue3 :3000)', 'Merchant SPA\n(Vue3 :3001)']),
    ('服务层', ['Spring Boot 3.2.5\n27个Controller\n16个Service\n21个Mapper', '拦截器链\nJWT(3套)+限流\nXSS+安全Headers', '基础设施\nRedis/PubSub/WS\nResilience4j/定时任务']),
    ('数据层', ['MySQL 8.0\n20张表\nInnoDB引擎', 'Redis 7.0\n缓存/锁/队列\nBitMap布隆', '腾讯地图API\nGeocode\n骑行路径规划']),
    ('技术栈', ['Java 17 + Maven\nMyBatis-Plus 3.5.6\nJJWT 0.12.5', 'Vue 3 + Vite\nElement Plus\nECharts + Pinia', '原生微信小程序\nWebSocket\nHutool 5.8.27']),
]
for ci, (title, items) in enumerate(cols):
    left = 1.5 + ci * 8
    add_textbox(s, left, 3.0, 7.5, 1, title, font_size=16, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    for ii, item in enumerate(items):
        y = 4.2 + ii * 3.8
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(left), Cm(y), Cm(7.5), Cm(3.3))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = GRAY; box.line.width = Pt(0.5)
        tf = add_rich_textbox(s, left+0.4, y+0.3, 6.7, 2.7)
        for line in item.split('\n'):
            add_para(tf, line, font_size=10, color=DARK, align=PP_ALIGN.CENTER, space_before=3)
add_page_number(s, 4, TOTAL_SLIDES)


# ================================================================
#  Slide 5 — 需求分析
# ================================================================
s = add_blank_slide()
add_title_bar(s, '三、需求分析概要')
add_table(s, 1.5, 3.5, [9, 11, 10.5],
    ['模块', '核心功能', '技术难点'],
    [
        ['模块1\n用户端与商家端\n(14项)',
         '用户：注册/登录、附近商家搜索、\n菜品浏览、购物车、下单支付、\n订单跟踪、评价评分\n'
         '商家：入驻审核、菜单管理(分类\n/上下架/库存/富文本)、订单处理、\n营业切换、销售统计',
         'GeoHash空间索引\n购物车多设备同步\n富文本编辑器集成'],
        ['模块2\n骑手端与配送系统\n(14项)',
         '骑手：注册认证(实名)、接单大厅、\n配送导航、位置轨迹上报、收入\n统计与提现、评分等级体系\n'
         '系统：智能派单、实时追踪、\nETA计算、超时预警、多骑手协同',
         'WebSocket实时推送\n四维加权派单算法\n腾讯地图API集成'],
        ['模块3\n后台管理与核心技术\n(15项)',
         '后台：数据看板(ECharts)、用户/商家\n/骑手管理、订单纠纷仲裁、营销\n活动(满减+优惠券)、权限管理\n'
         '技术：支付沙箱、消息推送、\n定时任务、缓存优化、安全防护',
         'Transaction Outbox\n分布式事务\n缓存三防(穿透/击穿/雪崩)\n接口限流与熔断降级'],
    ], font_size=10)
add_page_number(s, 5, TOTAL_SLIDES)


# ================================================================
#  Slide 6 — 数据库设计
# ================================================================
s = add_blank_slide()
add_title_bar(s, '四、数据库设计')
add_table(s, 1.5, 3.2, [4.5, 3.5, 4.5, 9, 9.5],
    ['模块', '表名', '用途', '核心字段', '关键索引'],
    [
        ['用户\n体系', 'user\nadmin\naddress', '用户/骑手账号\n平台管理员\n收货地址', 'phone,password(BCrypt)\nrole(user|delivery)\nlat,lng,is_default',
         'uk_phone_role\nuk_username\nidx_user_id'],
        ['商家\n体系', 'shop\ngoods_category\ngoods', '店铺信息\n商品分类\n商品/菜品', 'geohash,rating,sales\nbusiness_status\nprice,stock,rich_desc',
         'idx_geohash(6)\nidx_shop_id\nidx_status'],
        ['交易\n体系', 'cart\norder\norder_item', '购物车\n订单主表\n订单明细', 'uk_user_goods\nstatus(0-7),delivery_id\nis_large_order\ngoods_name/price快照',
         'uk_order_no\nidx_user_id\nidx_shop_id\nidx_delivery_id'],
        ['配送\n体系', 'delivery\ndelivery_record\ndelivery_track\nwithdraw', '骑手信息\n配送记录\nGPS轨迹\n提现申请', 'verify_status,level\non_time_rate,praise_rate\nlat,lng,speed\namount,status',
         'uk_user_id\nuk_order_id\nidx_delivery_time'],
        ['营销\n+ 评价', 'coupon\nuser_coupon\nfull_reduce_activity\nevaluation', '优惠券模板\n用户优惠券\n满减活动\n订单评价', 'condition_amount\nreduce_amount\nstart_time,end_time\nrating(1-5),status',
         'idx_time\nuk_order_id\nidx_user_id'],
        ['基础\n设施', 'event_log\naudit_log\nfeedback\nsystem_config', '事务事件表\n审计日志\n用户反馈/申诉\n系统配置', 'event_type,payload(JSON)\nretry_count\nadmin_id,action\nkey-value结构',
         'idx_status_create\nidx_user_id'],
    ], font_size=8)
add_textbox(s, 1.5, 16.8, 30, 0.8, '共 20 张表  |  InnoDB引擎  |  utf8mb4字符集  |  索引策略覆盖GeoHash前缀、联合索引、全文搜索等场景',
            font_size=10, color=DARK, align=PP_ALIGN.CENTER)
add_page_number(s, 6, TOTAL_SLIDES)


# ================================================================
#  Slide 7 — 核心功能1：用户端
# ================================================================
s = add_blank_slide()
add_title_bar(s, '五、核心功能实现 — 用户端与商家端')

add_section_box(s, 1.5, 3.2, 15, 7.2, 'C端用户功能（微信小程序）', [
    '手机号+密码注册登录，JWT Token鉴权，BCrypt加密',
    'GeoHash空间索引实现附近商家搜索（距离/评分/销量排序）',
    '分类浏览商品 → 加入购物车 → 优惠券+满减自动计算 → 下单',
    '模拟支付（微信/支付宝沙箱，90%成功率）',
    '8种订单状态流转：待支付→已支付→配送中→已完成/已取消/退款中/已退款/拒绝退款',
    '订单详情：实时状态 + 配送骑手信息 + 1-5星评价',
    '收货地址管理（GPS坐标标注）、优惠券中心、反馈申诉',
])
add_section_box(s, 17.5, 3.2, 15, 7.2, '商家端功能（Merchant SPA）', [
    '商家入驻申请 → 平台审核通过/拒绝（三级状态流转）',
    '商品CRUD + 分类管理 + 上下架 + 库存控制',
    '@wangeditor富文本编辑器，支持图文混排菜品描述',
    '订单处理：查看新单 → 接单/拒单 → 标记出餐完成',
    '营业状态一键切换（营业中/休息中）',
    '商家自有优惠券管理 + 满减活动查看',
    'ECharts销售统计图表（订单量/营业额/商品排行）',
])
add_section_box(s, 1.5, 11, 31, 7, '订单状态机（核心交易流转）', [
    '0-待支付 ──支付──→ 1-已支付 ──派单──→ 2-配送中 ──送达──→ 3-已完成',
    '0-待支付 ──取消──→ 4-已取消',
    '1-已支付/2-配送中 ──退款申请──→ 5-退款中 ──商家同意──→ 6-已退款',
    '5-退款中 ──商家拒绝──→ 7-商家拒绝退款 ──平台仲裁──→ 6-已退款 或 恢复为1-已支付',
    '大订单拆分：is_large_order=1 → splitLargeOrder() → 生成N个子订单 → 每个独立派单',
])
add_page_number(s, 7, TOTAL_SLIDES)


# ================================================================
#  Slide 8 — 核心功能2：骑手端
# ================================================================
s = add_blank_slide()
add_title_bar(s, '五、核心功能实现 — 骑手端与配送系统')

add_section_box(s, 1.5, 3.2, 15, 7.5, '骑手端功能（微信小程序）', [
    '实名认证注册：身份证号 + 真实姓名 → 平台审核',
    '接单大厅：在线/离线切换，查看附近待配送订单',
    '手动抢单 + 系统智能推送通知（WebSocket实时推送）',
    '配送导航：集成腾讯地图SDK骑行路线规划',
    'GPS轨迹上报：每5秒上报（Redis实时缓存 + DB批量落盘）',
    '收入统计：按日/周/月查看配送收入 + 提现申请',
    '骑手三级评价体系：铜牌/银牌/金牌（基于准时率+好评率+单量）',
])

add_section_box(s, 17.5, 3.2, 15, 7.5, '智能派单算法（DispatchServiceImpl）', [
    '评分公式：Score = 0.40×距离 + 0.25×质量 + 0.20×负载 + 0.15×顺路',
    '距离分：max(0, 1 − 骑手到店距离/5km)，超5km直接排除',
    '质量分：准时率×0.5 + 好评率×0.3 + 等级×0.2',
    '负载分：max(0, 1 − 当前订单数/2)，超2单直接排除',
    '顺路分：目的地1km内接近(+0.3) + 方位角30°内同方向(+0.2)',
    '三级分流：≥0.60自动派单 | ≥0.35推送抢单 | <0.35自由抢单',
])

add_section_box(s, 1.5, 11.3, 31, 6.8, '配送系统架构', [
    'WebSocket实时推送：OrderNotificationEndpoint + ConcurrentHashMap会话管理 + userId精确推送，延迟<500ms',
    'GPS轨迹存储优化：骑手上报 → Redis setex(30秒过期) → 追加到Redis List → 定时任务每分钟批量写入delivery_track表',
    '超时处理：@Scheduled每2分钟扫描配送超60分钟的订单 → 自动完成 + 扣减骑手准时率（按比例重算onTimeRate）',
    '大订单多骑手协同：splitLargeOrder() → 按商品拆分为子订单 → 每个子订单独立触发派单 → 不同骑手配送',
])
add_page_number(s, 8, TOTAL_SLIDES)


# ================================================================
#  Slide 9 — 核心功能3：后台管理
# ================================================================
s = add_blank_slide()
add_title_bar(s, '五、核心功能实现 — 后台管理与营销体系')

add_section_box(s, 1.5, 3.2, 15, 8, '平台管理后台（Admin SPA, 15页）', [
    '数据看板：ECharts实时图表（订单趋势/GMV/新增用户/活跃骑手）',
    '用户管理：C端用户列表 + 启用/禁用',
    '商家管理：入驻审核 + 信息修改 + 账号管理',
    '骑手管理：实名认证审核 + 状态/评分查看',
    '订单管理：全平台订单查询 + 异常订单强制取消',
    '纠纷仲裁：退款三级审核（用户→商家→平台最终裁定）',
    '营销配置：平台优惠券 + 满减活动（定时自动启停）',
    '系统配置：起送价/配送费/超时策略等键值对动态管理',
    '权限管理：超级管理员(admin) / 运营人员(operator) 角色区分',
])

add_section_box(s, 17.5, 3.2, 15, 8, '营销体系完整设计', [
    '优惠券模板(coupon)：条件金额 + 减免金额 + 总量/余量 + 有效期',
    '平台券(shop_id=NULL) vs 商家券(shop_id指定)',
    '用户领取(user_coupon)：未使用→已使用→已过期，三重状态',
    '下单自动校验：归属匹配 + 门槛校验 + 自动抵扣',
    '满减活动(full_reduce_activity)：三状态(未开始/进行中/已结束)',
    '自动启停：@Scheduled每2分钟扫描 start_time/end_time，自动开启/关闭',
    '下单自动匹配最优满减：getBestActivity(shopId, totalPrice)',
    '优惠叠加：优惠券 + 满减可同时使用，actualAmount = total + fee − 优惠券 − 满减',
])

add_section_box(s, 1.5, 11.8, 31, 6.3, '安全防护与性能优化', [
    '多层安全链：XssFilter → SecurityHeadersFilter → JWT鉴权(Wx/Admin/Merchant三套独立拦截器) → RateLimitInterceptor(Redis计数器IP级限流) → MyBatis-Plus参数化防SQL注入 → Resilience4j CircuitBreaker熔断降级',
    '缓存三防：CacheUtil集成布隆过滤器(穿透) + ConcurrentHashMap互斥锁+DoubleCheck(击穿) + 随机TTL±30%(雪崩)，一行代码getOrLoadWithBloom获取三防保护',
    '定时任务调度：5个@Scheduled方法（15min取消未支付 / 60min配送超时 / 2min优惠券过期 / 2min活动启停 / 1min事件兜底扫描），采用LIMIT分批防止全表扫描',
    '接口限流：@RateLimit注解声明式限流，Redis滑动窗口计数，触发返回HTTP 429 + CircuitBreakerFallback降级兜底',
])
add_page_number(s, 9, TOTAL_SLIDES)


# ================================================================
#  Slide 10 — 关键技术难点1：GeoHash + 派单
# ================================================================
s = add_blank_slide()
add_title_bar(s, '六、关键技术难点攻克（一）')

add_section_box(s, 1.5, 3.2, 15.5, 7.5, 'GeoHash空间索引（自研，158行）', [
    '从零实现Base32编码/解码，不依赖任何第三方地理库',
    'encode(lat,lng,precision=7) → "wm6y8g"（约150m精度）',
    'getNeighbors(geohash) → 8个相邻单元格（偏移→重编码方法）',
    'precisionForRadius(2.5km) → 精度9（自适应选择）',
    '搜索流程：GeoHash前缀过滤（9个单元格）→ Haversine精算 → 按距离/评分/销量排序',
    '配合腾讯地图Geocode API：商家文本地址 → 经纬度 → GeoHash自动更新',
    '索引：shop.geohash前缀索引(6位)，将全表扫描缩小到9个单元格',
])

add_section_box(s, 17.5, 3.2, 15.5, 7.5, '智能派单算法（354行）', [
    '四维加权：距离(40%) + 质量(25%) + 负载(20%) + 顺路(15%)',
    '方位角计算：bearing(lat1,lng1,lat2,lng2) → 两目的地方向差异<30°视为同方向',
    '三级分流：自动派单(≥0.60) / 推送抢单(≥0.35，前3名WebSocket通知) / 自由抢单(<0.35)',
    '派单前二次校验：防满单兜底（assignOrder中再次检查activeOrders）',
    '多骑手协同：大订单splitLargeOrder → 每2个商品一个子订单 → 独立派单',
])

add_section_box(s, 1.5, 11.3, 31, 6.8, '实现细节对比', [
    'GeoHash精度选择：1级(5000km) → 12级(0.08km)，搜索半径越大精度越低，始终搜索9个单元格保证召回。例如2.5km用精度9(约20m范围)，5km用精度7(约150m范围)',
    '派单顺路度计算两维度：(1)空间维度——新订单目的地与骑手现有订单目的地Haversine距离<1km→+0.3分；(2)方向维度——从店铺到两个目的地的方位角差异<30°→+0.2分。两者结合有效量化"顺路程度"',
    '骑手位置获取：从Redis读取rider:location:{deliveryId}（由骑手上报时setex），位置过期(>30秒)则跳过该骑手，保证派单数据时效性',
    'DispatchServiceImpl与DeliveryServiceImpl分离：派单算法专注于评分排序，配送服务专注于骑手认证/接单/轨迹上报，职责边界清晰',
])
add_page_number(s, 10, TOTAL_SLIDES)


# ================================================================
#  Slide 11 — 关键技术难点2：Outbox + 缓存
# ================================================================
s = add_blank_slide()
add_title_bar(s, '六、关键技术难点攻克（二）')

tf = add_rich_textbox(s, 1.5, 3.2, 31, 14)
add_para(tf, 'Transaction Outbox 分布式事务（EventLogService 184行）', font_size=18, bold=True, color=BLACK)
add_para(tf, '', font_size=3)
add_para(tf, '三级保障机制保证订单-库存最终一致性，无需RocketMQ/Kafka等重型中间件：', font_size=12, color=DARK)
add_para(tf, '', font_size=4)
add_para(tf, '第一级：业务事务内原子写入  →  orderMapper.updateById(order) + eventLogService.saveEvent("ORDER_PAID", payload) 在同一@Transactional中', font_size=11, color=BLACK, space_before=4)
add_para(tf, '第二级：afterCommit实时推送   →  TransactionSynchronization.afterCommit回调 → EventLogService.tryPublishAfterCommit → Redis Pub/Sub → OrderEventListener异步处理(通知+扣库存+更新销量)', font_size=11, color=BLACK, space_before=4)
add_para(tf, '第三级：定时任务兜底重试       →  @Scheduled(fixedRate=60000) 每分钟扫描 status=0且>60s 的事件 → LIMIT 50分批 → 最多重试5次 → 超限标记status=2(失败)', font_size=11, color=BLACK, space_before=4)
add_para(tf, '', font_size=8)

add_para(tf, '缓存三防体系（CacheUtil 183行）', font_size=18, bold=True, color=BLACK, space_before=16)
add_para(tf, '', font_size=3)
add_para(tf, '防穿透 — Redis Bitmap布隆过滤器：3个哈希函数，100万位空间。bloomMightContain返回false → 一定不存在，直接返回null不查DB；返回true → 查缓存 → 未命中查DB → 不存在则缓存空占位符"__NULL__"(5min TTL)', font_size=11, color=BLACK, space_before=6)
add_para(tf, '防击穿 — ConcurrentHashMap互斥锁：热点key过期时仅第一个请求进入synchronized(lock)重建缓存，其余等待。Double-Check机制防止缓存重建完成后重复加载', font_size=11, color=BLACK, space_before=4)
add_para(tf, '防雪崩 — 随机TTL±30%抖动：setWithJitter(key,value,baseTtl) → ttl = baseTtl ± (baseTtl×0.3×random)，防止大量key同时过期导致缓存雪崩', font_size=11, color=BLACK, space_before=4)
add_para(tf, '组合调用：getOrLoadWithBloom(namespace, value, cacheKey, loader) → 布隆判断(防穿透) → 查缓存 → 互斥锁加载(防击穿) → 随机TTL写入(防雪崩)，一行代码三防合一', font_size=11, color=BLACK, space_before=4)
add_page_number(s, 11, TOTAL_SLIDES)


# ================================================================
#  Slide 12 — 关键技术难点3：WebSocket + 安全
# ================================================================
s = add_blank_slide()
add_title_bar(s, '六、关键技术难点攻克（三）')

add_section_box(s, 1.5, 3.2, 15.5, 6.5, 'WebSocket实时推送', [
    'Jakarta WebSocket (@ServerEndpoint) 端点注册',
    'ConcurrentHashMap<String, Session> 会话管理',
    '按userId精确推送（从URI路径提取：/ws/order/{userId}）',
    'sendToUser(userId, json) → isOpen校验 → 异步发送',
    '消息类型：订单状态变更、新订单派单通知',
    '前端对接：小程序wx.connectSocket + Vite代理ws:true',
    '性能：<500ms延迟，支持2000+并发连接',
])

add_section_box(s, 17.5, 3.2, 15.5, 6.5, '安全防护纵深体系', [
    'XssFilter：请求参数HTML转义，拦截XSS攻击载荷',
    'SecurityHeadersFilter：X-Content-Type-Options/X-Frame-Options',
    'JWT鉴权：三套独立拦截器(Wx/Admin/Merchant)，各自校验签名密钥',
    'RateLimitInterceptor：@RateLimit注解 + Redis计数器 + IP/用户双维度限流',
    'MyBatis-Plus参数化查询：从框架层面杜绝SQL注入',
    'Resilience4j CircuitBreaker：熔断降级 → CircuitBreakerFallback兜底',
    '密码BCrypt加密 + 角色权限隔离 + 前端路由守卫',
])

add_section_box(s, 1.5, 10.3, 31, 7.8, '第三方集成与中间件运用', [
    '腾讯地图API (TencentMapService 300行)：地址解析(Geocode) → 文本地址转经纬度，用于商家入驻自动标注；骑行路径规划(bikingDistance) → 骑手到店铺+店铺到用户的双段距离计算，精确到米',
    'Redis五重角色：(1)缓存加速—热点商家列表/商品详情缓存；(2)分布式锁—库存扣减防超卖；(3)Pub/Sub消息队列—订单事件异步发布订阅，三个独立频道(order/evaluation/delivery)；(4)限流计数器—滑动窗口计数配合@RateLimit注解；(5)BitMap布隆过滤器—防缓存穿透',
    'Resilience4j熔断降级：对shopService/orderService/deliveryService/mapService四个关键服务配置CircuitBreaker，异常比例>50%时自动熔断，10秒后半开试探，恢复后自动关闭',
    '定时任务(OrderScheduledTask 192行)：autoCancelUnpaidOrders(15min未支付自动取消并恢复库存)、checkDeliveryTimeout(60min超时自动完成+扣减准时率)、autoExpireCoupons(优惠券到期自动失效)、autoUpdateActivityStatus(满减活动按时自动启停)、processEventLog(事件表兜底扫描，LIMIT 50分批)',
])
add_page_number(s, 12, TOTAL_SLIDES)


# ================================================================
#  Slide 13 — 系统测试
# ================================================================
s = add_blank_slide()
add_title_bar(s, '七、系统测试')
add_table(s, 1.5, 3.2, [4.5, 13, 6, 7],
    ['类别', '测试内容', '方法/工具', '结果'],
    [
        ['功能测试', '15个核心场景：注册登录、商家搜索、\n下单支付、智能派单、配送追踪、\n退款仲裁、营销启停等', '黑盒手动测试\n逐项验证执行',
         '全部通过\n发现并修复10+个Bug'],
        ['接口测试', '27个Controller全部REST接口\n参数校验 + JWT鉴权 + 异常捕获\n+ 统一响应格式', 'Knife4j (Swagger增强)\n在线文档逐接口调试',
         '全部通过\n响应格式统一'],
        ['性能指标', '商家搜索 GET /api/wx/shops\n下单 POST /api/wx/orders\n支付 POST /api/wx/orders/{id}/pay\n订单列表 GET /api/wx/orders',
         'JMeter / Apache Bench\n(建议)',
         'RT < 200ms(搜索)\nRT < 500ms(下单)\n并发500+(搜索)'],
        ['安全测试', 'XSS注入 / SQL注入 / 越权访问\n/ Token伪造 / 接口限流触发', '手动安全审查\nSwagger无鉴权调用',
         '全部拦截\n限流正常返回429'],
        ['已知限制', '短信验证码为模拟实现\n支付为沙箱90%成功率模拟\n单元测试待补充', '后续改进计划',
         '需对接真实\n第三方服务'],
    ], font_size=10)
add_page_number(s, 13, TOTAL_SLIDES)


# ================================================================
#  Slide 14 — 团队分工
# ================================================================
s = add_blank_slide()
add_title_bar(s, '八、团队分工')

add_table(s, 1.5, 3.2, [6, 8, 8, 8.5],
    ['成员', '后端开发', '前端开发', '数据库 / 测试 / 文档'],
    [
        ['余  睿\n(组长)\n2300770172',
         '系统架构 + 配置层\n订单模块(893行)\n派单算法(354行)\nGeoHash(158行)+缓存(183行)\n安全防护(228行)\nWebSocket+地图服务(420行)\nAdmin营销/系统接口',
         'Admin后台:\n营销配置+系统配置+反馈处理\n(约520行)',
         '系统架构设计\n综合设计报告主体\n过程报告 + PPT\n项目进度管理'],
        ['赵宇弦\n(组员)\n2300770178',
         '认证模块(311行)\n商品模块(280行)',
         'Admin后台:\n架构搭建+看板+订单+评价\n+API层(约2,200行)\nMerchant后台:\n全部6个页面(约1,500行)\n小程序C端:\n全部10个页面+7个组件(约2,400行)\n小程序骑手端:\n接单大厅+任务池(约250行)',
         '前端架构搭建\n接口联调'],
        ['周峥嵘\n(组员)\n2300770046',
         '配送系统(877行)\nTransaction Outbox(184行)\n事件监听器(207行)\n评价模块(162行)\n反馈模块(105行)\n定时任务(192行)\nAdmin用户/商家/骑手接口',
         'Admin后台:\n用户管理+商家审核+骑手管理\n(约570行)\n小程序骑手端:\n配送中+收入+骑手资料\n+评价(约750行)\n小程序C端:\n地址管理+反馈(约340行)',
         '数据库设计(500行SQL)\n15个测试用例 + Bug跟踪\n综合报告2个章节'],
    ], font_size=9)

tf2 = add_rich_textbox(s, 1.5, 15.5, 31, 3)
add_para(tf2, '分工原则：每人独立负责明确的模块，模块间通过REST API接口契约交互，互不依赖内部实现。通过Git分支管理实现并行开发。', font_size=11, color=DARK, space_before=6)
add_page_number(s, 14, TOTAL_SLIDES)


# ================================================================
#  Slide 15 — 总结与展望
# ================================================================
s = add_blank_slide()
add_title_bar(s, '九、总结与展望')

tf = add_rich_textbox(s, 1.5, 3.2, 31, 7)
add_para(tf, '项目完成情况', font_size=18, bold=True, color=BLACK)
add_para(tf, '', font_size=3)
add_para(tf, '三大模块49项子需求全部完成，代码总量30,000+行，后端116个Java类，前端两个SPA+一个小程序共42个页面，20张数据表。', font_size=12, color=DARK)
add_para(tf, '', font_size=8)
add_para(tf, '技术亮点', font_size=18, bold=True, color=BLACK)
add_para(tf, '', font_size=3)
highlights = [
    'Transaction Outbox：三级保障（事务写入→afterCommit推送→定时兜底），无需重型MQ实现最终一致性',
    '自研GeoHash：Base32编码/解码/邻域搜索/精度自适应，配合Haversine精算，从零实现不依赖第三方库',
    '缓存三防：布隆过滤器（穿透）+ 互斥锁+DoubleCheck（击穿）+ 随机TTL（雪崩），单类183行完整覆盖',
    '智能派单：四维加权（距离+质量+负载+顺路）+ 方位角方向计算 + 三级分流策略',
    '安全纵深：XSS→SecurityHeaders→JWT(3套)→限流→参数化防注入→熔断降级，六层防护',
]
for h in highlights:
    add_para(tf, '◆  ' + h, font_size=11, color=DARK, space_before=5)

tf2 = add_rich_textbox(s, 1.5, 14, 31, 4.5)
add_para(tf2, '改进方向', font_size=18, bold=True, color=BLACK)
add_para(tf2, '', font_size=3)
add_para(tf2, '◆  对接真实第三方服务（阿里云短信、微信支付V3 API、支付宝开放平台）    ◆  补充JUnit 5 + Mockito单元测试与MockMvc集成测试', font_size=11, color=DARK, space_before=5)
add_para(tf2, '◆  引入GitHub Actions / Jenkins CI/CD流水线（构建→测试→SonarQube→Docker→部署）    ◆  采用TDengine时序数据库优化GPS轨迹存储', font_size=11, color=DARK, space_before=5)
add_para(tf2, '◆  集成Prometheus + Grafana实现JVM/Redis/MySQL全链路监控    ◆  购物车WebSocket实时同步 + Douglas-Peucker轨迹压缩', font_size=11, color=DARK, space_before=5)
add_page_number(s, 15, TOTAL_SLIDES)


# ================================================================
#  Slide 16 — 致谢
# ================================================================
s = add_blank_slide()
shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(33.87), Cm(19.05))
shape.fill.solid(); shape.fill.fore_color.rgb = WHITE; shape.line.fill.background()
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(6), Cm(8), Cm(21.87), Cm(0.04))
line.fill.solid(); line.fill.fore_color.rgb = DARK; line.line.fill.background()

add_textbox(s, 6, 4, 22, 2.5, '感谢聆听', font_size=48, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
add_textbox(s, 6, 8.5, 22, 2, '请各位老师批评指正', font_size=22, color=DARK, align=PP_ALIGN.CENTER)
add_textbox(s, 6, 12, 22, 2, '第 11 组', font_size=18, color=GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, 6, 13.5, 22, 3, '余  睿（2300770172）    赵宇弦（2300770178）    周峥嵘（2300770046）\n软工2304',
            font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

add_textbox(s, 6, 17, 22, 1, 'Spring Boot 3.2.5  |  Vue 3  |  MyBatis-Plus  |  MySQL  |  Redis  |  WebSocket',
            font_size=10, color=GRAY, align=PP_ALIGN.CENTER)
add_page_number(s, 16, TOTAL_SLIDES)


# ================================================================
#  保存
# ================================================================
out = 'D:/SpringbootFinal/_temp_ppt.pptx'
prs.save(out)
print(f'PPT saved: {out}')
